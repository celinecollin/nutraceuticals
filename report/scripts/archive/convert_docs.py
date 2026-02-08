
import os
import glob
import pypandoc
import docx
from pptx import Presentation
import pandas as pd
import warnings

warnings.filterwarnings("ignore")

SOURCE_DIR = "."
OUTPUT_DIR = "report/sources"

def convert_docx_to_md(filepath, output_path):
    # Try pandoc first, fallback to python-docx
    try:
        output = pypandoc.convert_file(filepath, 'markdown', outputfile=output_path)
        print(f"Converted {filepath} to {output_path} (Pandoc)")
    except Exception as e:
        print(f"Pandoc failed for {filepath}, trying python-docx fallback...")
        try:
            doc = docx.Document(filepath)
            text = [f"# {os.path.basename(filepath)}"]
            for para in doc.paragraphs:
                # Basic formatting preservation
                prefix = ""
                if para.style.name.startswith('Heading'):
                    try:
                        level = int(para.style.name.split(' ')[-1])
                        prefix = "#" * level + " "
                    except:
                        prefix = "## "
                
                if para.style.name == 'List Bullet':
                    prefix = "- "
                
                text.append(f"{prefix}{para.text}")
            
            with open(output_path, "w") as f:
                f.write("\n\n".join(text))
            print(f"Converted {filepath} to {output_path} (python-docx)")
        except Exception as e2:
            print(f"Failed to convert {filepath}: {e2}")

def convert_pptx_to_md(filepath, output_path):
    try:
        prs = Presentation(filepath)
        md_content = f"# {os.path.basename(filepath)}\n\n"
        
        for i, slide in enumerate(prs.slides):
            md_content += f"## Slide {i+1}\n\n"
            if slide.shapes.title:
                md_content += f"### {slide.shapes.title.text}\n\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text != slide.shapes.title.text:
                    clean_text = shape.text.strip()
                    if clean_text:
                        md_content += f"{clean_text}\n\n"
        
        with open(output_path, "w") as f:
            f.write(md_content)
        print(f"Converted {filepath} to {output_path}")

    except Exception as e:
        print(f"Failed to convert {filepath}: {e}")

def convert_xlsx_to_md(filepath, output_path):
    try:
        # Simple sheet dump
        xl = pd.ExcelFile(filepath)
        with open(output_path, "w") as f:
            f.write(f"# {os.path.basename(filepath)}\n\n")
            for sheet in xl.sheet_names:
                f.write(f"## Sheet: {sheet}\n\n")
                try:
                    df = xl.parse(sheet)
                    # Convert to markdown using tabulate (pandas uses it internally for markdown)
                    if not df.empty:
                        f.write(df.to_markdown(index=False))
                    else:
                        f.write("(Empty Sheet)")
                except Exception as e:
                    f.write(f"(Error reading sheet: {e})")
                f.write("\n\n")
        print(f"Converted {filepath} to {output_path}")
    except Exception as e:
        print(f"Failed to convert {filepath}: {e}")

def main():
    if not os.path.exists(OUTPUT_DIR):
        os.makedirs(OUTPUT_DIR)

    # Process all files in root
    for ext, converter in [
        ('docx', convert_docx_to_md),
        ('pptx', convert_pptx_to_md),
        ('xlsx', convert_xlsx_to_md)
    ]:
        # Case insensitive glob
        files = []
        files.extend(glob.glob(f"*.{ext}"))
        files.extend(glob.glob(f"*.{ext.upper()}"))
        
        # Unique files
        files = list(set(files))
        
        print(f"Found {len(files)} {ext} files.")
        
        for fp in files:
            # Skip temp files
            if os.path.basename(fp).startswith('~$'):
                continue
                
            name = os.path.basename(fp)
            output_name = f"{os.path.splitext(name)[0]}.md"
            output_path = os.path.join(OUTPUT_DIR, output_name)
            converter(fp, output_path)

if __name__ == "__main__":
    main()
